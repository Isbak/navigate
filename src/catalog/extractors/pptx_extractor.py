from __future__ import annotations

from collections.abc import Iterable
from pathlib import Path
from typing import Any

from pptx import Presentation


def _iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    """Yield shapes recursively, including shapes inside groups."""

    for shape in shapes:
        yield shape
        grouped_shapes = getattr(shape, "shapes", None)
        if grouped_shapes is not None:
            yield from _iter_shapes(grouped_shapes)


def _hyperlink_address(shape: Any) -> str | None:
    """Return a shape hyperlink address when the shape type supports one."""

    try:
        click_action = getattr(shape, "click_action", None)
    except TypeError:
        # python-pptx raises for group shapes because groups cannot have their
        # own click action. Their child shapes are inspected separately.
        return None
    if click_action is None:
        return None
    hyperlink = getattr(click_action, "hyperlink", None)
    if hyperlink is None:
        return None
    return getattr(hyperlink, "address", None)


class PptxExtractor:
    def extract_text(self, path: Path) -> str:
        prs = Presentation(path)
        parts: list[str] = []
        for slide in prs.slides:
            for shape in _iter_shapes(slide.shapes):
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
                if address := _hyperlink_address(shape):
                    parts.append(address)
        return "\n".join(parts)
