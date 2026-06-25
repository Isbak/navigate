"""Generate deterministic binary test fixtures for the extractor comparison study.

Run once (checked-in output is in benchmarks/corpus/binary_fixtures/):
    python -m benchmarks.make_binary_fixtures
"""

from __future__ import annotations

from pathlib import Path

FIXTURE_DIR = Path(__file__).parent / "corpus" / "binary_fixtures"


# ---------------------------------------------------------------------------
# PDF fixtures
# ---------------------------------------------------------------------------

def _make_simple_text_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text(
        (72, 100),
        "Introduction to Navigate\n\n"
        "Navigate is a local-first knowledge platform that catalogs\n"
        "documents without moving source files.\n\n"
        "Key Features:\n"
        "- Fast offline extraction\n"
        "- Hyperlink discovery\n"
        "- Knowledge graph construction\n\n"
        "Visit https://github.com/isbak/navigate for more information.",
        fontsize=12,
    )
    doc.save(str(path))
    doc.close()


def _make_multi_column_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    # Left column
    page.insert_text(
        (50, 100),
        "Section A\n\nThis is the left column.\nIt describes the first topic.\n\n"
        "Architecture\nThe system uses a modular\nextraction pipeline.",
        fontsize=10,
    )
    # Right column
    page.insert_text(
        (320, 100),
        "Section B\n\nThis is the right column.\nIt describes the second topic.\n\n"
        "Performance\nExtraction runs in parallel\nusing thread workers.",
        fontsize=10,
    )
    doc.save(str(path))
    doc.close()


def _make_table_heavy_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text(
        (72, 80),
        "Extraction Backend Comparison\n\n"
        "Backend   | Format | Quality | Speed\n"
        "----------|--------|---------|------\n"
        "Current   | PDF    | Fair    | Fast\n"
        "MarkItDown| DOCX   | Good    | Fast\n"
        "Docling   | PDF    | Best    | Slow\n\n"
        "See https://docling.example.com for Docling documentation.",
        fontsize=11,
    )
    doc.save(str(path))
    doc.close()


def _make_hyperlink_pdf(path: Path) -> None:
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text(
        (72, 100),
        "Reference Documents\n\n"
        "The architecture specification is at https://docs.example.com/architecture\n"
        "The API reference is at https://api.example.com/reference\n"
        "The deployment guide is at https://ops.example.com/deploy\n\n"
        "For questions, contact the platform team.",
        fontsize=12,
    )
    doc.save(str(path))
    doc.close()


# ---------------------------------------------------------------------------
# DOCX fixtures
# ---------------------------------------------------------------------------

def _make_docx_with_table(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Design Specification", level=1)
    doc.add_paragraph(
        "This document describes the architecture for the catalog service."
    )
    doc.add_heading("Component Overview", level=2)

    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = table.rows[0].cells
    headers[0].text = "Component"
    headers[1].text = "Responsibility"
    headers[2].text = "Owner"
    data = [
        ("Scanner", "Discovers and indexes files", "Platform Team"),
        ("Extractor", "Extracts text and links", "Platform Team"),
        ("Classifier", "LLM-based knowledge extraction", "AI Team"),
    ]
    for row, (comp, resp, owner) in zip(table.rows[1:], data, strict=False):
        row.cells[0].text = comp
        row.cells[1].text = resp
        row.cells[2].text = owner

    doc.add_paragraph(
        "\nFor more information, see the architecture guide at "
        "https://confluence.example.com/architecture"
    )

    # Add a hyperlink relationship (embedded, not inline text)
    part = doc.part
    r_id = part.relate_to(
        "https://jira.example.com/ADR-001",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    _ = r_id  # relationship stored in the OPC package

    doc.save(str(path))


def _make_docx_simple(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Release Governance Policy", level=1)
    doc.add_paragraph(
        "All production releases must pass the release governance checklist "
        "before deployment. The checklist is reviewed by the Architecture Board."
    )
    doc.add_heading("Approval Process", level=2)
    doc.add_paragraph(
        "1. Submit the release candidate to the release queue.\n"
        "2. The Architecture Decision Record must be approved.\n"
        "3. Security review must be completed.\n"
        "4. Performance benchmarks must meet the SLA targets."
    )
    doc.add_paragraph(
        "See https://confluence.example.com/release-governance for the full policy."
    )
    doc.save(str(path))


# ---------------------------------------------------------------------------
# PPTX fixtures
# ---------------------------------------------------------------------------

def _make_pptx_slides(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # title + content

    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.title.text = "Platform Architecture"
    slide1.placeholders[1].text = (
        "The Navigate platform consists of four main components:\n"
        "• Scanner: file system discovery\n"
        "• Extractor: text and link extraction\n"
        "• Classifier: LLM-based knowledge graph construction\n"
        "• API: REST interface for downstream consumers"
    )

    slide2 = prs.slides.add_slide(blank_layout)
    slide2.shapes.title.text = "Extraction Pipeline"
    slide2.placeholders[1].text = (
        "The extraction pipeline supports multiple modes:\n"
        "fast — text-only, no API calls\n"
        "enhanced — MarkItDown for office formats\n"
        "docling — IBM Docling for superior PDF quality\n"
        "high-quality — Claude vision for equations and scans\n\n"
        "See https://docs.example.com/extraction for details."
    )

    prs.save(str(path))


# ---------------------------------------------------------------------------
# XLSX fixtures
# ---------------------------------------------------------------------------

def _make_xlsx_workbook(path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = "Metrics"

    headers = ["Metric", "Baseline", "Target", "Status"]
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = Font(bold=True)

    data = [
        ("Extraction throughput (docs/s)", 12.5, 15.0, "In Progress"),
        ("Link discovery precision", 0.94, 0.95, "Met"),
        ("Classification recall", 0.82, 0.85, "In Progress"),
        ("Knowledge graph coverage", 0.71, 0.80, "Needs Work"),
    ]
    for row_idx, row_data in enumerate(data, 2):
        for col_idx, value in enumerate(row_data, 1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    # Add hyperlinks to cells
    ws["E1"] = "Reference"
    ws["E1"].font = Font(bold=True)
    ws["E2"].hyperlink = "https://grafana.example.com/d/perf-dashboard"
    ws["E2"].value = "Dashboard"
    ws["E3"].hyperlink = "https://confluence.example.com/sla-targets"
    ws["E3"].value = "SLA Doc"

    wb.save(str(path))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    FIXTURE_DIR.mkdir(parents=True, exist_ok=True)

    makers = [
        ("simple_text.pdf", _make_simple_text_pdf),
        ("multi_column.pdf", _make_multi_column_pdf),
        ("table_heavy.pdf", _make_table_heavy_pdf),
        ("hyperlink.pdf", _make_hyperlink_pdf),
        ("table_docx.docx", _make_docx_with_table),
        ("simple_docx.docx", _make_docx_simple),
        ("slide_deck.pptx", _make_pptx_slides),
        ("data_workbook.xlsx", _make_xlsx_workbook),
    ]

    for filename, maker in makers:
        dest = FIXTURE_DIR / filename
        maker(dest)
        print(f"  created {dest.relative_to(Path(__file__).parent.parent)}")

    print(f"\n{len(makers)} fixtures written to {FIXTURE_DIR}")


if __name__ == "__main__":
    main()
